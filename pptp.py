from pptx import Presentation
from pptx.util import Inches

# Create a presentation object
prs = Presentation()

# Slide 1: Title Slide
slide_1 = prs.slides.add_slide(prs.slide_layouts[0])
title = slide_1.shapes.title
subtitle = slide_1.placeholders[1]

title.text = "CO2 Emission Prediction Project"
subtitle.text = "Using Machine Learning and Data Visualization to Predict Future Emissions\nPresented by Luis Llamas"

# Slide 2: Project Overview
slide_2 = prs.slides.add_slide(prs.slide_layouts[1])
title = slide_2.shapes.title
title.text = "Project Overview"
content = slide_2.shapes.placeholders[1].text = """
- The CO2 Emission Prediction Project aims to analyze historical CO2 emission data and predict future emissions.
- The project uses a combination of data science, machine learning, and web technologies.
- It includes an interactive web-based dashboard for visualization and predictions based on historical data.
"""

# Slide 3: Problem Statement
slide_3 = prs.slides.add_slide(prs.slide_layouts[1])
title = slide_3.shapes.title
title.text = "Problem Statement"
content = slide_3.shapes.placeholders[1].text = """
- Climate change is driven by rising CO2 emissions.
- Tools for forecasting future emissions can help governments and organizations in decision-making.
- This project builds a predictive model for CO2 emissions and presents results via a web-based dashboard.
"""

# Slide 4: Project Workflow
slide_4 = prs.slides.add_slide(prs.slide_layouts[1])
title = slide_4.shapes.title
title.text = "Project Workflow"
content = slide_4.shapes.placeholders[1].text = """
1. Data Collection: Gather historical CO2 emissions data.
2. Data Cleaning & Preprocessing: Handle missing values and normalize data.
3. Data Analysis & Visualization: Explore data patterns.
4. Machine Learning Model: Train a Linear Regression model.
5. Web Development: Build an interactive dashboard using Flask and Plotly.
"""

# Slide 5: Technologies Used
slide_5 = prs.slides.add_slide(prs.slide_layouts[1])
title = slide_5.shapes.title
title.text = "Technologies Used"
content = slide_5.shapes.placeholders[1].text = """
- Data Processing: Python, Pandas, Jupyter Notebooks
- Machine Learning: Scikit-learn (Linear Regression)
- Web Development: Flask, Plotly, HTML/CSS, Bootstrap
- Database: PostgreSQL, SQLAlchemy
- Cloud & Deployment: AWS Elastic Beanstalk, AWS RDS
"""

# Slide 6: Data Overview
slide_6 = prs.slides.add_slide(prs.slide_layouts[1])
title = slide_6.shapes.title
title.text = "Data Overview"
content = slide_6.shapes.placeholders[1].text = """
- Historical CO2 emissions data by country (1990-2020).
- Data Source: Global environmental datasets.
- Key features: CO2 emissions per country (in kilotons), Years 1990 to 2020.
"""

# Slide 7: Data Preprocessing
slide_7 = prs.slides.add_slide(prs.slide_layouts[1])
title = slide_7.shapes.title
title.text = "Data Preprocessing"
content = slide_7.shapes.placeholders[1].text = """
- Removed missing values using imputation.
- Re-formatted columns and handled null values.
- Normalized data for consistency across years and countries.
"""

# Slide 8: Machine Learning Model
slide_8 = prs.slides.add_slide(prs.slide_layouts[1])
title = slide_8.shapes.title
title.text = "Machine Learning Model"
content = slide_8.shapes.placeholders[1].text = """
- Model: Linear Regression
- Training Data: Historical data from 1990-2020
- Evaluation Metrics: Mean Squared Error (MSE) and R-squared
- Predicted future CO2 emissions for 2025-2027.
"""

# Slide 9: Web Application
slide_9 = prs.slides.add_slide(prs.slide_layouts[1])
title = slide_9.shapes.title
title.text = "Web Application Overview"
content = slide_9.shapes.placeholders[1].text = """
- Flask-based web app with interactive visualizations.
- Data visualization of historical CO2 emissions.
- Prediction tool for future CO2 emissions (2025-2027).
- Frontend: HTML/CSS/Bootstrap and Plotly for interactive charts.
"""

# Slide 10: Data Visualizations
slide_10 = prs.slides.add_slide(prs.slide_layouts[1])
title = slide_10.shapes.title
title.text = "Data Visualizations"
content = slide_10.shapes.placeholders[1].text = """
- Line chart: Global CO2 emissions over time.
- Bar chart: Top 10 CO2 emitting countries (2020).
- Line chart: CO2 emissions over time for selected countries.
- Prediction charts for future emissions (2025-2027).
"""

# Slide 11: Deployment
slide_11 = prs.slides.add_slide(prs.slide_layouts[1])
title = slide_11.shapes.title
title.text = "Cloud Deployment on AWS"
content = slide_11.shapes.placeholders[1].text = """
- Deployed using AWS Elastic Beanstalk.
- PostgreSQL database managed on AWS RDS.
- Flask app runs on Elastic Beanstalk with an EC2 instance.
- Code managed with Git and version control.
"""

# Slide 12: Project Demo
slide_12 = prs.slides.add_slide(prs.slide_layouts[1])
title = slide_12.shapes.title
title.text = "Live Demo"
content = slide_12.shapes.placeholders[1].text = "Present a live demo of the web app."

# Slide 13: Future Improvements
slide_13 = prs.slides.add_slide(prs.slide_layouts[1])
title = slide_13.shapes.title
title.text = "Future Improvements"
content = slide_13.shapes.placeholders[1].text = """
- Use more complex models like Random Forest or XGBoost.
- Integrate real-time data from external APIs.
- Expand the model to include other environmental factors (e.g., deforestation, renewable energy).
"""

# Slide 14: Challenges Faced
slide_14 = prs.slides.add_slide(prs.slide_layouts[1])
title = slide_14.shapes.title
title.text = "Challenges Faced"
content = slide_14.shapes.placeholders[1].text = """
- Data Cleaning: Handling inconsistencies and missing values.
- Model Tuning: Optimizing the machine learning model for accuracy.
- Deployment: Managing Python version compatibility issues on AWS.
"""

# Slide 15: Conclusion
slide_15 = prs.slides.add_slide(prs.slide_layouts[1])
title = slide_15.shapes.title
title.text = "Conclusion"
content = slide_15.shapes.placeholders[1].text = """
- The CO2 Emission Prediction Project demonstrates the use of machine learning for environmental data analysis.
- Predictive tools like this can help policymakers make informed decisions to combat climate change.
"""

# Slide 16: Thank You
slide_16 = prs.slides.add_slide(prs.slide_layouts[1])
title = slide_16.shapes.title
title.text = "Thank You!"
content = slide_16.shapes.placeholders[1].text = """
For more information, please contact:
Luis Llamas
Email: luis.llamas@maslla.com
"""

# Save the presentation
prs.save('CO2_Emission_Prediction_Project_Presentation.pptx')
print("Presentation created successfully!")
